"""Generate CV Website Presentation PPTX — Visual redesign."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw
import tempfile, os

# ── Colours ───────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0D, 0x14)
CARD_BG  = RGBColor(0x13, 0x12, 0x1F)
DEEP     = RGBColor(0x15, 0x14, 0x22)
PURPLE   = RGBColor(0x86, 0x7B, 0xFF)
PURPLE2  = RGBColor(0x63, 0x50, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0xAA, 0xAA, 0xBB)
DIMGREY  = RGBColor(0x55, 0x55, 0x70)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
GOLD     = RGBColor(0xFB, 0xBF, 0x24)
INDIGO   = RGBColor(0xA5, 0xB4, 0xFC)
RED      = RGBColor(0xF8, 0x71, 0x71)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
CODE_BG  = RGBColor(0x0A, 0x0A, 0x12)
BORDER   = RGBColor(0x28, 0x24, 0x48)
DOT_R    = RGBColor(0xFF, 0x5F, 0x57)
DOT_Y    = RGBColor(0xFF, 0xBD, 0x2E)
DOT_G    = RGBColor(0x28, 0xCA, 0x41)
BLUE_LI  = RGBColor(0x0A, 0x66, 0xC2)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Temp file registry (cleanup at end) ───────────────────────────────────
_tmp_files = []

# ── Image helpers ─────────────────────────────────────────────────────────
def make_circle_png(img_path, px=600):
    img = Image.open(img_path).convert("RGBA")
    sq  = min(img.size)
    lc  = (img.width - sq) // 2
    tc  = (img.height - sq) // 2
    img = img.crop((lc, tc, lc + sq, tc + sq)).resize((px, px), Image.LANCZOS)

    mask = Image.new('L', (px, px), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, px - 1, px - 1), fill=255)
    result = Image.new('RGBA', (px, px), (0, 0, 0, 0))
    result.paste(img, mask=mask)

    # Purple border ring
    bw = 12
    ImageDraw.Draw(result).ellipse(
        (bw // 2, bw // 2, px - bw // 2, px - bw // 2),
        outline=(0x86, 0x7B, 0xFF, 255), width=bw
    )
    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    result.save(tmp.name, 'PNG')
    tmp.close()
    _tmp_files.append(tmp.name)
    return tmp.name


CIRCLE_PROFILE = make_circle_png(r"c:\Websites\MyCV\profile.jpg")
SCREENSHOT_PNG  = r"c:\Websites\MyCV\image.png"

# ── Core helpers ──────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def rect(slide, l, t, w, h, fill=CARD_BG, border=BORDER, bw=0.75, radius=False):
    shape = slide.shapes.add_shape(5 if radius else 1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(bw)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, cx, cy, r, fill, border=None):
    shape = slide.shapes.add_shape(9, cx - r, cy - r, r * 2, r * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def hline(slide, l, t, w, color=PURPLE, thickness=1.5):
    shape = slide.shapes.add_shape(1, l, t, w, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def vline(slide, l, t, h, color=BORDER, thickness=1):
    shape = slide.shapes.add_shape(1, l, t, Pt(thickness), h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def tb(slide, text, l, t, w, h, size=14, bold=False,
       color=WHITE, align=PP_ALIGN.LEFT, italic=False, font="Arial"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text         = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return box


def bullets(slide, items, l, t, w, h, size=13, color=GREY, gap=Pt(5)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = gap
        run = p.add_run()
        indent = item.startswith("  ")
        run.text = ("  › " if indent else "• ") + item.lstrip()
        run.font.size      = Pt(size - (1 if indent else 0))
        run.font.color.rgb = DIMGREY if indent else color
        run.font.name      = "Arial"


def pill_shape(slide, text, l, t, w, h, bg, fg):
    s = slide.shapes.add_shape(5, l, t, w, h)   # rounded rect
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.fill.background()
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top  = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = fg; r.font.name = "Arial"


def section_bar(slide, title, t=Inches(0.32)):
    """Purple-accented section title."""
    tb(slide, title, Inches(0.55), t, Inches(12.2), Inches(0.52),
       size=26, bold=True, color=WHITE)
    hline(slide, Inches(0.55), t + Inches(0.56), Inches(12.2))


def code_block(slide, code, l, t, w, h, font_size=9.2):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = CODE_BG
    s.line.color.rgb = BORDER; s.line.width = Pt(0.75)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top  = tf.margin_bottom = Inches(0.1)
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run(); r.text = line
        r.font.size = Pt(font_size); r.font.name = "Courier New"
        r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)


def browser(slide, code, l, t, w, h, url="", font_size=9.2):
    """Code block inside a browser chrome frame."""
    bar_h = Inches(0.3)
    # Outer window
    rect(slide, l, t, w, h, fill=CODE_BG, border=BORDER, bw=1)
    # Title bar
    rect(slide, l, t, w, bar_h, fill=DEEP, border=None)
    # 3 dots
    ds = Inches(0.082)
    for i, dc in enumerate([DOT_R, DOT_Y, DOT_G]):
        circle(slide, l + Inches(0.14) + i * Inches(0.15),
               t + bar_h / 2, ds / 2, dc)
    # URL bar
    uw = w - Inches(0.66)
    ul = l + Inches(0.55)
    rect(slide, ul, t + Inches(0.07), uw, Inches(0.17), fill=CODE_BG, border=BORDER, bw=0.5)
    if url:
        tb(slide, url, ul + Inches(0.06), t + Inches(0.073), uw - Inches(0.1), Inches(0.15),
           size=6.5, color=DIMGREY)
    # Code area
    code_block(slide, code, l, t + bar_h, w, h - bar_h, font_size=font_size)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()

# Background gradient-like bars
for i in range(6):
    alpha = int(0x18 - i * 2)
    rect(sl, 0, Inches(i * 1.25), W, Inches(1.3),
         fill=RGBColor(0x0D, 0x0D, max(0x14, 0x20 - i * 2)), border=None)

# Left purple top bar
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)

# Profile photo (circular) – right side
prof_w = Inches(2.7)
sl.shapes.add_picture(CIRCLE_PROFILE,
                       W - prof_w - Inches(0.65),
                       (H - prof_w) / 2,
                       prof_w, prof_w)

# Title text – left side
tb(sl, "Personal CV Website",
   Inches(0.55), Inches(1.55), Inches(9.2), Inches(1.3),
   size=50, bold=True, color=WHITE)

hline(sl, Inches(0.55), Inches(3.0), Inches(5.5))

tb(sl, "Design, Features & Code Walkthrough",
   Inches(0.55), Inches(3.1), Inches(9.2), Inches(0.55),
   size=20, color=PURPLE)

tb(sl, "Htet Myat Aung",
   Inches(0.55), Inches(3.75), Inches(9.2), Inches(0.45),
   size=17, bold=True, color=WHITE)

tb(sl, "BSc Computer Science (AI) · Royal Holloway, University of London",
   Inches(0.55), Inches(4.2), Inches(9.2), Inches(0.38),
   size=13, color=GREY)

# Tech stack pills at bottom
stack = ["HTML", "CSS", "JavaScript", "5 Pages", "Dark Theme"]
px = Inches(0.55)
for s_text in stack:
    pw = Inches(len(s_text) * 0.13 + 0.6)
    pill_shape(sl, s_text, px, Inches(5.0), pw, Inches(0.32),
               bg=RGBColor(0x1E, 0x1A, 0x44), fg=PURPLE)
    px += pw + Inches(0.12)

# Bottom tag line
tb(sl, "10-minute presentation",
   Inches(0.55), Inches(6.8), Inches(6), Inches(0.35),
   size=11, color=DIMGREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – PROJECT OVERVIEW (visual stat cards)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "Project Overview")

# Description
tb(sl, "A hand-coded personal portfolio & CV website — no frameworks, no templates.",
   Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.4),
   size=15, color=GREY)

# 4 large stat cards
stats = [
    ("5", "HTML Pages", PURPLE, RGBColor(0x1E,0x1A,0x44)),
    ("7", "Projects Showcased", GREEN, RGBColor(0x0D,0x2A,0x22)),
    ("18", "Countries Reached", GOLD, RGBColor(0x2A,0x22,0x08)),
    ("1st", "Class Average", INDIGO, RGBColor(0x18,0x18,0x38)),
]
card_w = Inches(2.9)
card_h = Inches(1.55)
gap    = Inches(0.25)
top_t  = Inches(1.7)
lx     = Inches(0.55)
for val, label, fg, bg in stats:
    rect(sl, lx, top_t, card_w, card_h, fill=bg, border=fg, bw=1)
    tb(sl, val, lx, top_t + Inches(0.15), card_w, Inches(0.75),
       size=44, bold=True, color=fg, align=PP_ALIGN.CENTER)
    tb(sl, label, lx, top_t + Inches(0.85), card_w, Inches(0.45),
       size=12, color=WHITE, align=PP_ALIGN.CENTER)
    lx += card_w + gap

# What it includes
tb(sl, "What the website includes",
   Inches(0.55), Inches(3.5), Inches(12.2), Inches(0.38),
   size=14, bold=True, color=PURPLE)

features = [
    ("Home", "Hero · Stats · Featured projects · Skills · Contact"),
    ("About", "Extended bio · Education · Quick-facts sidebar · CTA"),
    ("Projects", "7 project cards · Filter tabs · MyanPay video feature"),
    ("Experience", "3 roles · DICE Lab research · achievement bullets"),
    ("Contact", "Form with JS validation · Direct links · Availability"),
]
fx = Inches(0.55)
for title, desc in features:
    fw = Inches(2.4)
    rect(sl, fx, Inches(3.95), fw, Inches(1.05), fill=CARD_BG, border=BORDER)
    tb(sl, title, fx + Inches(0.12), Inches(4.05), fw - Inches(0.2), Inches(0.32),
       size=12, bold=True, color=PURPLE)
    tb(sl, desc, fx + Inches(0.12), Inches(4.37), fw - Inches(0.2), Inches(0.55),
       size=9.5, color=GREY)
    fx += fw + Inches(0.18)

# Single shared stylesheet note
hline(sl, Inches(0.55), Inches(5.18), Inches(12.2), color=BORDER)
tb(sl, "One shared  main.css  (1,100+ lines)  ·  Per-page ES module JavaScript  ·  Responsive (600px & 700px breakpoints)",
   Inches(0.55), Inches(5.28), Inches(12.2), Inches(0.4),
   size=12, color=DIMGREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – DESIGN & PREPARATION PROCESS (visual timeline)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "Design & Preparation Process")

steps = [
    ("01", "Planning",        PURPLE,  "Mapped 5 pages & content. Identified audience:\nrecruiters & researchers."),
    ("02", "Design System",   GREEN,   "Dark theme #0d0d14. Purple #867bff accent.\nConsistent colour tokens across all pages."),
    ("03", "Typography",      GOLD,    "Google Fonts: Syne (headings) + DM Mono\n(labels). Loaded via <link preconnect>."),
    ("04", "Layout Strategy", INDIGO,  "CSS Grid for multi-column sections.\nFlexbox for nav & inline groups."),
    ("05", "JS Architecture", ORANGE,  "animations.js shared module — imported by\nall 5 page JS files. Keeps code DRY."),
    ("06", "Testing",         RED,     "Chrome + Firefox. DevTools responsive mode\nfor 600px & 700px breakpoints."),
]

# Horizontal timeline spine
spine_y = Inches(3.5)
spine_l = Inches(0.8)
spine_w = Inches(11.7)
rect(sl, spine_l, spine_y - Pt(1), spine_w, Pt(2), fill=BORDER, border=None)

step_w  = spine_w / len(steps)
for i, (num, title, color, desc) in enumerate(steps):
    cx = spine_l + step_w * i + step_w / 2

    # Circle node on spine
    r = Inches(0.22)
    circle(sl, cx, spine_y, r, fill=color)
    tb(sl, num, cx - r, spine_y - r, r * 2, r * 2,
       size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

    # Connector line going up or down alternately
    above = (i % 2 == 0)
    card_h2 = Inches(1.1)
    gap2 = Inches(0.45)
    card_t = spine_y - r - gap2 - card_h2 if above else spine_y + r + gap2
    line_t = (spine_y - r) if above else (spine_y + r)
    line_h = gap2
    rect(sl, cx - Pt(0.75), line_t if above else spine_y + r,
         Pt(1.5), gap2, fill=color, border=None)

    # Info card
    card_l = cx - Inches(0.85)
    rect(sl, card_l, card_t, Inches(1.7), card_h2,
         fill=CARD_BG, border=color, bw=1)
    tb(sl, title, card_l + Inches(0.1), card_t + Inches(0.08),
       Inches(1.5), Inches(0.3), size=11, bold=True, color=color)
    tb(sl, desc, card_l + Inches(0.1), card_t + Inches(0.36),
       Inches(1.5), Inches(0.7), size=8.5, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – WEBSITE ARCHITECTURE (visual diagram)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "Website Architecture")

# ── HTML pages row ──
html_pages = ["index.html", "about.html", "projects.html", "experience.html", "contact.html"]
html_colors = [PURPLE, GREEN, GOLD, INDIGO, ORANGE]
page_w = Inches(2.1)
page_h = Inches(0.75)
page_gap = Inches(0.22)
page_top = Inches(1.2)
page_lx  = (W - (len(html_pages) * page_w + (len(html_pages)-1) * page_gap)) / 2

for i, (name, col) in enumerate(zip(html_pages, html_colors)):
    lx = page_lx + i * (page_w + page_gap)
    rect(sl, lx, page_top, page_w, page_h, fill=CARD_BG, border=col, bw=1.5)
    tb(sl, name, lx, page_top, page_w, page_h,
       size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

# Label
tb(sl, "5 HTML Pages", Inches(0.3), page_top + Inches(0.18),
   Inches(1.8), Inches(0.38), size=11, color=DIMGREY, bold=True)

# Shared CSS node – centre
css_w = Inches(2.6)
css_h = Inches(0.65)
css_l = (W - css_w) / 2
css_t = Inches(2.45)
rect(sl, css_l, css_t, css_w, css_h, fill=RGBColor(0x1E,0x1A,0x44), border=PURPLE, bw=2)
tb(sl, "main.css", css_l, css_t, css_w, css_h,
   size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
tb(sl, "Shared by all 5 pages", Inches(0.3), css_t + Inches(0.12),
   Inches(1.8), Inches(0.38), size=11, color=DIMGREY, bold=True)

# Lines from each HTML page down to CSS
css_cx = css_l + css_w / 2
css_ty = css_t
for i, col in enumerate(html_colors):
    px_cx = page_lx + i * (page_w + page_gap) + page_w / 2
    rect(sl, px_cx - Pt(0.75), page_top + page_h,
         Pt(1.5), css_ty - (page_top + page_h), fill=col, border=None)

# JS files row
js_files = [
    ("animations.js", PURPLE, "shared module"),
    ("main.js",       PURPLE, "home"),
    ("about.js",      GREEN,  "about"),
    ("projects.js",   GOLD,   "projects"),
    ("experience.js", INDIGO, "experience"),
    ("contact.js",    ORANGE, "contact"),
]
js_w   = Inches(1.8)
js_h   = Inches(0.7)
js_gap = Inches(0.14)
js_top = Inches(3.6)
js_lx  = (W - (len(js_files) * js_w + (len(js_files)-1) * js_gap)) / 2
for i, (name, col, sub) in enumerate(js_files):
    lx = js_lx + i * (js_w + js_gap)
    rect(sl, lx, js_top, js_w, js_h, fill=CARD_BG, border=col, bw=1.5)
    tb(sl, name, lx, js_top + Inches(0.05), js_w, Inches(0.36),
       size=9.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, sub, lx, js_top + Inches(0.38), js_w, Inches(0.26),
       size=8, color=DIMGREY, align=PP_ALIGN.CENTER)

# "animations.js is imported by all" arrow / label
anim_cx = js_lx + js_w / 2
anim_by = js_top + js_h
rect(sl, anim_cx - Pt(1), anim_by, Pt(2), Inches(0.45), fill=PURPLE, border=None)
tb(sl, "exported & imported →", js_lx + js_w + Inches(0.1), anim_by + Inches(0.08),
   Inches(3.5), Inches(0.3), size=9, color=DIMGREY, italic=True)
tb(sl, "JS Labels", Inches(0.3), js_top + Inches(0.18),
   Inches(1.8), Inches(0.38), size=11, color=DIMGREY, bold=True)

# Key facts below
hline(sl, Inches(0.55), Inches(4.78), Inches(12.2), color=BORDER)
facts = [
    "No build tools — runs directly in the browser",
    "type=\"module\" enables ES import/export between JS files",
    "Google Fonts loaded with <link preconnect> for zero layout shift",
    "Single CSS file = one source of truth for all design tokens",
]
bullets(sl, facts, Inches(0.55), Inches(4.9), Inches(12.2), Inches(1.6),
        size=12, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – HTML: NAVIGATION
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "HTML — Navigation & Shared Header")

nav_code = """\
<div class="header">
  <a href="index.html">
    <img src="profile.jpg" class="nav-logo">
  </a>

  <a class="nav-link active-page" href="about.html">About</a>
  <a class="nav-link" href="projects.html">Projects</a>
  <a class="nav-link" href="experience.html">Experience</a>
  <a class="nav-link" href="index.html#skills">Skills</a>
  <a class="nav-link" href="contact.html">Contact</a>

  <!-- Dropdown for secondary links -->
  <div class="nav-dropdown">
    <button class="nav-dropdown-toggle">More &#9662;</button>
    <div class="dropdown-menu">
      <a href="https://github.com/HtetMyatAungg"
         target="_blank" rel="noopener noreferrer">GitHub</a>
      <a href="...drive.google.com/...">Download CV</a>
      <a href="mailto:henry.htetmyataung@gmail.com">Email Me</a>
    </div>
  </div>

  <button id="hire-me">LinkedIn</button>
</div>"""

browser(sl, nav_code, Inches(0.45), Inches(1.12), Inches(6.35), Inches(5.9),
        url="index.html — navigation", font_size=9.5)

# Mini visual nav mockup on the right
nav_box_l = Inches(7.05)
nav_box_t = Inches(1.12)
nav_box_w = Inches(5.95)
nav_box_h = Inches(0.55)
rect(sl, nav_box_l, nav_box_t, nav_box_w, nav_box_h,
     fill=RGBColor(0x0D,0x0D,0x14), border=BORDER)
# Logo circle
circle(sl, nav_box_l + Inches(0.22), nav_box_t + nav_box_h/2,
       Inches(0.14), PURPLE)
# Nav items
nav_labels = ["About", "Projects", "Experience", "Skills", "Contact", "More ▾"]
link_x = nav_box_l + Inches(0.52)
for j, lbl in enumerate(nav_labels):
    col = WHITE if j == 0 else DIMGREY
    tb(sl, lbl, link_x, nav_box_t + Inches(0.12), Inches(0.9), Inches(0.3),
       size=8.5, color=col, align=PP_ALIGN.CENTER)
    link_x += Inches(0.9)
# LinkedIn button
rect(sl, nav_box_l + nav_box_w - Inches(0.9), nav_box_t + Inches(0.1),
     Inches(0.78), Inches(0.35), fill=PURPLE2, border=None)
tb(sl, "LinkedIn", nav_box_l + nav_box_w - Inches(0.88),
   nav_box_t + Inches(0.12), Inches(0.76), Inches(0.28),
   size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "↑ what the browser renders", nav_box_l, nav_box_t + nav_box_h + Inches(0.04),
   nav_box_w, Inches(0.25), size=8.5, color=DIMGREY, italic=True)

# Explanation bullets
bullets(sl, [
    "position:sticky; top:0 — navbar stays fixed while page scrolls",
    "active-page class = white text on current page, grey on others",
    "Profile <img> doubles as home button (click logo → index.html)",
    "Dropdown menu uses CSS :hover — no JavaScript needed",
    "LinkedIn button (#hire-me) styled & linked in each page's JS file",
    "rel=\"noopener noreferrer\" on external links — security best practice",
    "Smooth anchor scroll: JS catches href=\"#skills\" → scrollIntoView()",
], Inches(7.05), Inches(1.88), Inches(5.95), Inches(4.8),
   size=12, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – CSS: DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "CSS — Design System & Colour Palette")

# Colour swatches (large visual)
swatches = [
    (BG,      "#0d0d14", "Background"),
    (CARD_BG, "#13121f", "Card Surface"),
    (PURPLE,  "#867bff", "Accent & CTA"),
    (INDIGO,  "#a5b4fc", "AI/ML Badge"),
    (GREEN,   "#34d399", "Live Badge"),
    (GOLD,    "#fbbf24", "Hackathon"),
    (RED,     "#f87171", "Form Error"),
    (BLUE_LI, "#0a66c2", "LinkedIn"),
]
sw_w = Inches(1.43)
sw_h = Inches(0.75)
sw_gap = Inches(0.16)
sw_l = Inches(0.55)
sw_t = Inches(1.18)
for i, (rgb, hex_v, lbl) in enumerate(swatches):
    lx = sw_l + i * (sw_w + sw_gap)
    rect(sl, lx, sw_t, sw_w, sw_h, fill=rgb, border=BORDER, bw=0.5)
    tb(sl, hex_v, lx, sw_t + sw_h + Inches(0.03), sw_w, Inches(0.22),
       size=8.5, color=GREY, align=PP_ALIGN.CENTER, font="Courier New")
    tb(sl, lbl, lx, sw_t + sw_h + Inches(0.25), sw_w, Inches(0.22),
       size=8, color=DIMGREY, align=PP_ALIGN.CENTER)

# Left: Typography specimen
tb(sl, "Typography", Inches(0.55), Inches(2.68), Inches(5.8), Inches(0.35),
   size=13, bold=True, color=PURPLE)
rect(sl, Inches(0.55), Inches(3.05), Inches(5.8), Inches(3.8),
     fill=CARD_BG, border=BORDER)
tb(sl, "Syne", Inches(0.75), Inches(3.2), Inches(5.4), Inches(0.55),
   size=34, bold=True, color=WHITE, font="Arial")
tb(sl, "SKILL CATEGORY LABELS — WEIGHT 700",
   Inches(0.75), Inches(3.75), Inches(5.4), Inches(0.35),
   size=10, bold=True, color=PURPLE, font="Arial")
hline(sl, Inches(0.75), Inches(4.15), Inches(5.2), color=BORDER)
tb(sl, "DM Mono", Inches(0.75), Inches(4.25), Inches(5.4), Inches(0.5),
   size=28, color=WHITE, font="Courier New")
tb(sl, "skill-item labels · metric values · code",
   Inches(0.75), Inches(4.75), Inches(5.4), Inches(0.32),
   size=10, color=GREY, font="Courier New")
hline(sl, Inches(0.75), Inches(5.1), Inches(5.2), color=BORDER)
tb(sl, "Arial / Helvetica — body text fallback",
   Inches(0.75), Inches(5.2), Inches(5.4), Inches(0.32),
   size=11, color=DIMGREY, font="Arial")
tb(sl, "Loaded via Google Fonts <link preconnect>",
   Inches(0.75), Inches(5.52), Inches(5.4), Inches(0.28),
   size=9, color=DIMGREY, italic=True)

# Right: CSS code for global reset + title rule
tb(sl, "Global CSS Reset", Inches(6.7), Inches(2.68), Inches(6.3), Inches(0.35),
   size=13, bold=True, color=PURPLE)
global_css = """\
/* Dark theme applied to EVERY element */
* {
  color: white;
  background-color: #0d0d14;
  font-family: Arial, Helvetica, sans-serif;
  padding-right: 1em;
  padding-left: 1em;
}

/* Responsive media — never overflow */
img, video, iframe {
  max-width: 100%;
  height: auto;
}

/* Section title with purple rule after text */
.title {
  display: flex;
  align-items: center;
  gap: 8px;
}
.title::after {
  content: "";
  flex: 1;
  height: 1px;
  background: rgb(72, 36, 156);
}

/* Accent text + pill for badge */
.intro-text {
  color: #867bff;
  background-color: #857bf82d;
  border-radius: 1em;
  padding: 1em;
}"""
browser(sl, global_css, Inches(6.7), Inches(3.05), Inches(6.3), Inches(3.8),
        url="main.css", font_size=9.2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – CSS: LAYOUT (Grid + Flexbox visual diagrams)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "CSS — Grid & Flexbox Layouts")

# === Left side: visual diagram of project grid ===
tb(sl, "CSS Grid — Projects Page", Inches(0.55), Inches(1.12),
   Inches(5.8), Inches(0.35), size=13, bold=True, color=PURPLE)

# Grid demo
grid_l = Inches(0.55)
grid_t = Inches(1.55)
grid_w = Inches(5.8)
cell_w = (grid_w - Inches(0.2)) / 3
cell_h = Inches(1.15)
cell_gap = Inches(0.1)
card_labels = ["MyanPay\n(Web)", "MyanLearn\n(Web)", "NASA NEO\n(ML)",
               "Face Clf\n(ML)", "CF Review\n(Web)", "Minesweeper\n(ML)"]
card_cols = [PURPLE, GREEN, INDIGO, INDIGO, GREEN, INDIGO]
for ci, (lbl, col) in enumerate(zip(card_labels, card_cols)):
    col_i = ci % 3
    row_i = ci // 3
    cx = grid_l + col_i * (cell_w + cell_gap)
    cy = grid_t + row_i * (cell_h + cell_gap)
    rect(sl, cx, cy, cell_w, cell_h, fill=CARD_BG, border=col, bw=1)
    tb(sl, lbl, cx, cy + Inches(0.35), cell_w, Inches(0.5),
       size=9, color=col, align=PP_ALIGN.CENTER)
tb(sl, "grid-template-columns: repeat(auto-fill, minmax(300px, 1fr))",
   grid_l, grid_t + 2*(cell_h + cell_gap) + Inches(0.08),
   grid_w, Inches(0.28), size=8.5, color=DIMGREY, font="Courier New")

# Flexbox nav demo
tb(sl, "Flexbox — Navigation Bar", Inches(0.55), Inches(4.35),
   Inches(5.8), Inches(0.35), size=13, bold=True, color=PURPLE)
nav_demo_t = Inches(4.75)
rect(sl, Inches(0.55), nav_demo_t, Inches(5.8), Inches(0.55),
     fill=DEEP, border=BORDER)
items = ["⬤", "About", "Projects", "Exp", "Skills", "Contact", "More", "LinkedIn"]
item_w = Inches(0.68)
for ji, itm in enumerate(items):
    col_it = PURPLE if ji == 0 else (WHITE if ji == 1 else DIMGREY)
    tb(sl, itm, Inches(0.6) + ji * item_w, nav_demo_t + Inches(0.1),
       item_w, Inches(0.32), size=8, color=col_it, align=PP_ALIGN.CENTER)
tb(sl, "display:flex  justify-content:space-evenly  position:sticky  top:0",
   Inches(0.55), nav_demo_t + Inches(0.6), Inches(5.8), Inches(0.25),
   size=8.5, color=DIMGREY, font="Courier New")

# Media query demo
tb(sl, "Responsive Breakpoints", Inches(0.55), Inches(5.55),
   Inches(5.8), Inches(0.35), size=13, bold=True, color=PURPLE)
media_css = """\
@media (max-width: 600px) {
  .project-grid-new { grid-template-columns: 1fr; }
  .skills-grid       { grid-template-columns: 1fr; }
  .video-feature     { grid-template-columns: 1fr; }
  h1.name            { font-size: 1.8em; }
}"""
browser(sl, media_css, Inches(0.55), Inches(5.92), Inches(5.8), Inches(1.38),
        url="main.css — @media", font_size=9)

# === Right side: all CSS layout code ===
layout_css = """\
/* ── Projects: responsive auto-fill grid ── */
.project-grid-new {
  display: grid;
  grid-template-columns:
    repeat(auto-fill, minmax(300px, 1fr));
  gap: 1.25em;
}

/* ── About page: sidebar + bio ── */
.about-page-grid {
  display: grid;
  grid-template-columns: 260px 1fr;
  gap: 3em;
  align-items: start;
}

/* ── Skills: 3 equal columns ── */
.skills-grid {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 16px;
}

/* ── MyanPay feature: text | video ── */
.video-feature {
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 2em;
  align-items: center;
}

/* ── Sticky navigation ── */
.header {
  display: flex;
  justify-content: space-evenly;
  align-items: center;
  position: sticky;
  top: 0;
  z-index: 1000;
  border-bottom: 1px solid #131e7c;
}

/* ── CTA button row ── */
.buttons {
  display: flex;
  gap: 1em;
}"""
browser(sl, layout_css, Inches(6.7), Inches(1.12), Inches(6.3), Inches(6.18),
        url="main.css — layout", font_size=9.5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – CSS: CARDS, BADGES & ANIMATIONS
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "CSS — Cards, Badges & Scroll Animations")

# Left: card code
card_css = """\
/* ── Base project card ── */
.pcard {
  background: #13121f;
  border: 1px solid rgba(255,255,255,0.07);
  border-radius: 1em;
  padding: 1.4em 1.5em 1.2em;
  display: flex;
  flex-direction: column;
  gap: 0.6em;
  transition: transform 0.2s,
              border-color 0.2s,
              box-shadow 0.2s;
}
.pcard:hover {
  transform: translateY(-3px);      /* subtle lift */
  border-color: rgba(134,123,255,0.4);
  box-shadow: 0 8px 30px rgba(99,102,241,0.12);
}
.pcard.featured {
  background: linear-gradient(
    135deg, #13121f 0%, #1a1630 100%);
}

/* ── Scroll animation ── */
.animate-on-scroll {
  opacity: 0;
  transform: translateY(24px);
  transition: opacity 0.55s ease,
              transform 0.55s ease;
}
.animate-on-scroll.visible {
  opacity: 1;
  transform: translateY(0);
}"""
browser(sl, card_css, Inches(0.45), Inches(1.12), Inches(5.7), Inches(5.3),
        url="main.css — .pcard + .animate-on-scroll", font_size=9.3)

# Right top: visual card mock
tb(sl, "Card visual (hover → lifts 3px)", Inches(6.5), Inches(1.12),
   Inches(6.5), Inches(0.33), size=12, bold=True, color=PURPLE)

# Normal card
rect(sl, Inches(6.5), Inches(1.5), Inches(2.8), Inches(1.3),
     fill=CARD_BG, border=BORDER, bw=1)
tb(sl, "Normal", Inches(6.6), Inches(1.55), Inches(2.6), Inches(0.28),
   size=9, color=DIMGREY)
tb(sl, "Project Card", Inches(6.6), Inches(1.83), Inches(2.6), Inches(0.3),
   size=13, bold=True, color=WHITE)
tb(sl, "Some metric here", Inches(6.6), Inches(2.13), Inches(2.6), Inches(0.25),
   size=9, color=PURPLE)

# Hover card (slightly up, glowing border)
rect(sl, Inches(9.65), Inches(1.44), Inches(2.8), Inches(1.3),
     fill=CARD_BG, border=PURPLE, bw=1.5)
tb(sl, "On Hover", Inches(9.75), Inches(1.49), Inches(2.6), Inches(0.28),
   size=9, color=PURPLE)
tb(sl, "Project Card", Inches(9.75), Inches(1.77), Inches(2.6), Inches(0.3),
   size=13, bold=True, color=WHITE)
tb(sl, "Some metric here", Inches(9.75), Inches(2.07), Inches(2.6), Inches(0.25),
   size=9, color=PURPLE)
tb(sl, "↑ translateY(-3px)", Inches(9.65), Inches(2.78), Inches(2.8), Inches(0.25),
   size=9, italic=True, color=DIMGREY, align=PP_ALIGN.CENTER)

# Badge specimens
tb(sl, "Status badges", Inches(6.5), Inches(3.1), Inches(6.5), Inches(0.3),
   size=12, bold=True, color=PURPLE)
badge_data = [
    ("LIVE", GREEN, RGBColor(0x0D,0x2A,0x22)),
    ("HEDERA HACKATHON", GOLD, RGBColor(0x2A,0x22,0x08)),
    ("AI / ML", INDIGO, RGBColor(0x18,0x18,0x38)),
    ("IN PROGRESS", ORANGE, RGBColor(0x2A,0x18,0x08)),
]
bx = Inches(6.5)
for lbl, fg, bg in badge_data:
    bw2 = Inches(len(lbl) * 0.1 + 0.5)
    pill_shape(sl, lbl, bx, Inches(3.48), bw2, Inches(0.3), bg=bg, fg=fg)
    bx += bw2 + Inches(0.15)

# Animate-on-scroll visual
tb(sl, "Scroll animation — before & after", Inches(6.5), Inches(4.0),
   Inches(6.5), Inches(0.3), size=12, bold=True, color=PURPLE)
# Before (invisible + shifted)
for xi, (lbl, op) in enumerate([("opacity: 0\ntranslatY(24px)", DIMGREY),
                                  ("opacity: 1\ntranslatY(0)", GREEN)]):
    bx2 = Inches(6.5) + xi * Inches(3.1)
    rect(sl, bx2, Inches(4.38), Inches(2.7), Inches(0.85),
         fill=CARD_BG, border=op, bw=1.5)
    tb(sl, ["Before (hidden)", "After .visible"][xi],
       bx2 + Inches(0.1), Inches(4.45), Inches(2.5), Inches(0.25),
       size=9, bold=True, color=op)
    tb(sl, lbl, bx2 + Inches(0.1), Inches(4.7), Inches(2.5), Inches(0.45),
       size=9, color=GREY, font="Courier New")

# Explanation bullets at bottom
bullets(sl, [
    "translateY(-3px) hover lift: subtle depth cue keeping the dark theme clean",
    "RGBA colour badges use transparency — they blend into the card surface naturally",
    "opacity:0 + translateY(24px) = invisible until JS fires — prevents flash of content",
], Inches(0.45), Inches(6.5), Inches(12.5), Inches(1.0),
   size=11.5, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 – JS: SCROLL ANIMATIONS
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "JavaScript — Scroll Animations  (animations.js)")

anim_code = """\
// Shared ES module — imported by every page
export function initScrollAnimations() {

  // All elements marked for animation
  const targets =
    document.querySelectorAll('.animate-on-scroll')
  if (!targets.length) return

  // Apply data-delay as CSS transition-delay
  // so cards stagger in one by one
  targets.forEach(el => {
    const delay = el.dataset.delay
    if (delay) el.style.transitionDelay = delay
  })

  // IntersectionObserver fires when 10% of element
  // enters the viewport — much better than scroll event
  const observer = new IntersectionObserver(
    (entries) => {
      entries.forEach(entry => {
        if (entry.isIntersecting) {
          // Add .visible → CSS transition plays
          entry.target.classList.add('visible')
          // Unobserve: animation fires exactly once
          observer.unobserve(entry.target)
        }
      })
    },
    { threshold: 0.1 }
  )

  targets.forEach(el => observer.observe(el))
}

// Usage in main.js (and every other page JS):
import { initScrollAnimations } from './animations.js'
initScrollAnimations()"""

browser(sl, anim_code, Inches(0.45), Inches(1.12), Inches(6.5), Inches(6.18),
        url="animations.js", font_size=9.5)

# Right: flow diagram
tb(sl, "How IntersectionObserver works",
   Inches(7.2), Inches(1.12), Inches(5.8), Inches(0.35),
   size=13, bold=True, color=PURPLE)

flow = [
    (PURPLE,  "Page loads\nElements hidden\n(opacity:0)"),
    (INDIGO,  "Observer\nwatches each\n.animate-on-scroll"),
    (GREEN,   "User scrolls\n10% of element\nenters viewport"),
    (GOLD,    "JS adds\n.visible class\nto element"),
    (WHITE,   "CSS transition\nplays: fade up\ninto view"),
]
box_w = Inches(1.0)
box_h = Inches(1.05)
box_gap = Inches(0.14)
flow_l = Inches(7.25)
flow_t = Inches(1.58)
for fi, (col, lbl) in enumerate(flow):
    fx = flow_l + fi * (box_w + box_gap)
    rect(sl, fx, flow_t, box_w, box_h, fill=CARD_BG, border=col, bw=1.5)
    tb(sl, lbl, fx, flow_t + Inches(0.08), box_w, box_h - Inches(0.1),
       size=8.5, color=col, align=PP_ALIGN.CENTER)
    if fi < len(flow) - 1:
        tb(sl, "→", fx + box_w, flow_t + box_h/2 - Inches(0.15),
           box_gap, Inches(0.3), size=12, color=DIMGREY, align=PP_ALIGN.CENTER)

# data-delay stagger visual
tb(sl, "Stagger effect with data-delay",
   Inches(7.2), Inches(2.82), Inches(5.8), Inches(0.32),
   size=12, bold=True, color=PURPLE)
stagger_labels = [
    ("data-delay=\"0s\"",    Inches(0.0),  PURPLE),
    ("data-delay=\"0.1s\"",  Inches(0.1),  GREEN),
    ("data-delay=\"0.2s\"",  Inches(0.2),  GOLD),
]
for si, (lbl, delay, col) in enumerate(stagger_labels):
    bar_l = Inches(7.25)
    bar_t = Inches(3.2) + si * Inches(0.55)
    bar_w_full = Inches(5.6)
    bar_h2 = Inches(0.38)
    rect(sl, bar_l, bar_t, bar_w_full, bar_h2, fill=CARD_BG, border=BORDER, bw=0.5)
    # fill showing delay
    fill_w = bar_w_full * (0.45 + si * 0.2)
    rect(sl, bar_l, bar_t, fill_w, bar_h2, fill=col, border=None)
    tb(sl, lbl, bar_l + Inches(0.1), bar_t + Inches(0.08),
       Inches(3), Inches(0.22), size=9, color=BG, bold=True)

bullets(sl, [
    "IntersectionObserver is more performant than a scroll event listener — browser-native, zero JS overhead",
    "threshold:0.1 triggers when 10% enters view — natural 'just coming in' timing",
    "observer.unobserve() after trigger — animation fires exactly once, no re-play",
    "CSS does the visual work; JS only toggles one class name (.visible)",
], Inches(7.2), Inches(4.95), Inches(5.9), Inches(2.3),
   size=11.5, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 – JS: PROJECT FILTER TABS
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "JavaScript — Project Filter Tabs  (projects.js)")

filter_code = """\
// ── Filter tabs ──
document.querySelectorAll(".filter-btn").forEach(btn => {
  btn.addEventListener("click", function() {

    // 1. Remove .active from all buttons
    document.querySelectorAll(".filter-btn")
      .forEach(b => b.classList.remove("active"))

    // 2. Mark clicked button as active
    btn.classList.add("active")

    // 3. Read filter from data-filter attribute
    const filter = btn.dataset.filter

    // 4. Show / hide cards by data-category
    document.querySelectorAll(".pcard").forEach(card => {
      if (filter === "all") {
        card.classList.remove("hidden")
      } else {
        const cats = card.dataset.category || ""
        card.classList.toggle(
          "hidden",
          !cats.split(" ").includes(filter)
        )
      }
    })
  })
})

// ── Video placeholder fallback ──
const video = document.getElementById("myanpay-video")
const placeholder = document.getElementById("video-placeholder")

if (video && placeholder) {
  // Start hidden — only reveal if video proves it loaded
  video.style.display = "none"
  placeholder.style.display = "flex"

  video.addEventListener("loadeddata", () => {
    placeholder.style.display = "none"
    video.style.display = "block"
  })
  video.addEventListener("error", () => {
    video.style.display = "none"
    placeholder.style.display = "flex"
  })
}"""

browser(sl, filter_code, Inches(0.45), Inches(1.12), Inches(6.5), Inches(6.18),
        url="projects.js", font_size=9.2)

# Right: visual filter mockup
tb(sl, "Filter tabs in action", Inches(7.2), Inches(1.12),
   Inches(5.8), Inches(0.32), size=13, bold=True, color=PURPLE)

filter_labels = [("All", True, PURPLE), ("Web & Edge", False, None),
                 ("AI / ML", False, None), ("In Progress", False, None)]
fx2 = Inches(7.25)
for lbl, active, _ in filter_labels:
    fw2 = Inches(len(lbl) * 0.1 + 0.6)
    bg2 = PURPLE if active else CARD_BG
    fg2 = WHITE  if active else DIMGREY
    pill_shape(sl, lbl, fx2, Inches(1.55), fw2, Inches(0.32), bg=bg2, fg=fg2)
    fx2 += fw2 + Inches(0.1)

# Grid showing which cards are visible / hidden
grid2_l = Inches(7.25); grid2_t = Inches(2.0)
mini_w = Inches(1.7); mini_h = Inches(0.85); mini_gap = Inches(0.12)
mini_cards = [
    ("MyanPay",    "web",    GREEN,  False),
    ("MyanLearn",  "web",    GREEN,  False),
    ("NASA NEO",   "ml",     INDIGO, False),
    ("LLM Scratch","ml wip", ORANGE, False),
    ("Face Clf",   "ml",     INDIGO, False),
    ("CF Reviewer","web",    GREEN,  False),
    ("Minesweeper","ml",     INDIGO, False),
]
for ci2, (nm, cat, col2, hidden) in enumerate(mini_cards):
    col_i = ci2 % 3; row_i = ci2 // 3
    mx = grid2_l + col_i * (mini_w + mini_gap)
    my = grid2_t  + row_i * (mini_h + mini_gap)
    border_col = col2 if not hidden else BORDER
    rect(sl, mx, my, mini_w, mini_h, fill=CARD_BG, border=border_col, bw=1)
    tb(sl, nm, mx, my + Inches(0.08), mini_w, Inches(0.3),
       size=9, bold=True, color=col2 if not hidden else DIMGREY,
       align=PP_ALIGN.CENTER)
    tb(sl, cat, mx, my + Inches(0.38), mini_w, Inches(0.25),
       size=7.5, color=DIMGREY, align=PP_ALIGN.CENTER)

bullets(sl, [
    "HTML: data-filter on buttons, data-category on cards (space-separated for multi-match)",
    "cats.split(' ').includes(filter) — card matches if filter is in its category list",
    "classList.toggle('hidden', condition) — one-liner to add or remove hidden",
    "CSS: .pcard.hidden { display: none } — hiding done in CSS, JS manages only the class",
    "Video fallback: shows placeholder by default, reveals video only on successful load",
], Inches(7.2), Inches(5.35), Inches(5.9), Inches(2.15),
   size=11.5, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 – JS: CONTACT FORM VALIDATION
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "JavaScript — Form Validation  (contact.js)")

form_code = """\
function showError(fieldId, message) {
  const el    = document.getElementById(fieldId+"-error")
  const input = document.getElementById(fieldId)
  if (el)    el.textContent = message
  if (input) input.classList.add("input-error")
}

function clearErrors() {
  document.querySelectorAll(".form-error")
    .forEach(el => el.textContent = "")
  document.querySelectorAll(".input-error")
    .forEach(el => el.classList.remove("input-error"))
  successMsg.style.display = "none"
}

form.addEventListener("submit", function(e) {
  e.preventDefault()   // stop browser default submit
  clearErrors()

  const name    = document.getElementById("fname").value.trim()
  const email   = document.getElementById("femail").value.trim()
  const subject = document.getElementById("fsubject").value
  const message = document.getElementById("fmessage").value.trim()
  let valid = true

  if (!name) {
    showError("fname", "Name is required.")
    valid = false
  }
  // Regex: must have @ and dot in domain
  if (!/^[^\\s@]+@[^\\s@]+\\.[^\\s@]+$/.test(email)) {
    showError("femail", "Enter a valid email.")
    valid = false
  }
  if (!subject) {
    showError("fsubject", "Please select a subject.")
    valid = false
  }
  if (message.length < 10) {
    showError("fmessage", "Min 10 characters.")
    valid = false
  }
  if (valid) {
    successMsg.style.display = "block"
    form.reset()
  }
})"""

browser(sl, form_code, Inches(0.45), Inches(1.12), Inches(6.4), Inches(6.18),
        url="contact.js", font_size=9.2)

# Right: visual form mockup
tb(sl, "Contact form — live validation",
   Inches(7.1), Inches(1.12), Inches(5.9), Inches(0.32),
   size=13, bold=True, color=PURPLE)

form_t = Inches(1.55)
form_l = Inches(7.15)
form_w = Inches(5.9)
fields = [
    ("Name", "Your name", False),
    ("Email", "your@email.com", True),
    ("Subject", "Select a subject ▾", False),
    ("Message", "Tell me what you are working on...", False),
]
for fi2, (lbl, ph, error) in enumerate(fields):
    fy = form_t + fi2 * Inches(0.82)
    tb(sl, lbl, form_l, fy, form_w, Inches(0.24),
       size=9, bold=True, color=GREY)
    field_h = Inches(0.35) if lbl != "Message" else Inches(0.55)
    border_c = RED if error else BORDER
    rect(sl, form_l, fy + Inches(0.25), form_w, field_h,
         fill=CARD_BG, border=border_c, bw=1)
    tb(sl, ph if not error else "Please enter a valid email address.",
       form_l + Inches(0.1), fy + Inches(0.28),
       form_w - Inches(0.15), field_h - Inches(0.08),
       size=8.5, color=RED if error else DIMGREY)
    if error:
        tb(sl, "✕ Please enter a valid email address.",
           form_l, fy + Inches(0.65), form_w, Inches(0.22),
           size=8, color=RED)

# Send button
btn_t = form_t + 4 * Inches(0.82) + Inches(0.2)
rect(sl, form_l, btn_t, Inches(1.6), Inches(0.38),
     fill=PURPLE, border=None)
tb(sl, "Send Message", form_l, btn_t + Inches(0.05), Inches(1.6), Inches(0.28),
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Success state
rect(sl, form_l, btn_t + Inches(0.48), form_w, Inches(0.38),
     fill=RGBColor(0x0D,0x2A,0x22), border=GREEN, bw=0.8)
tb(sl, "✓  Thanks! I will get back to you within 24 hours.",
   form_l + Inches(0.12), btn_t + Inches(0.55), form_w - Inches(0.2), Inches(0.28),
   size=9.5, color=GREEN)

bullets(sl, [
    "novalidate on <form> — disables browser default validation, we control all styles",
    "Regex /^[^\\s@]+@[^\\s@]+\\.[^\\s@]+$/ validates email without any library",
    "input-error CSS class turns field border red (#f87171) — instant visual feedback",
    "Real-time clear: errors remove as user starts typing — not frustrating to use",
], form_l, btn_t + Inches(1.05), form_w, Inches(1.5),
   size=11, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 – EXPERIENCE PAGE + RESEARCH SCREENSHOT
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "Experience Page — DICE Lab Research")

# Add actual screenshot
sl.shapes.add_picture(SCREENSHOT_PNG,
                       Inches(0.45), Inches(1.12),
                       Inches(7.5), Inches(4.2))
tb(sl, "↑ Autoformalization Research Website (dicelab-rhul.github.io) — built on the codebase I contributed to",
   Inches(0.45), Inches(5.38), Inches(7.5), Inches(0.35),
   size=9, color=DIMGREY, italic=True)

# Right: DICE Lab experience bullets
tb(sl, "DICE Lab · Volunteer Research Assistant",
   Inches(8.3), Inches(1.12), Inches(4.8), Inches(0.35),
   size=13, bold=True, color=PURPLE)
tb(sl, "Royal Holloway, University of London · Mar 2026–Present",
   Inches(8.3), Inches(1.5), Inches(4.8), Inches(0.28),
   size=10, color=DIMGREY)

bullets(sl, [
    "Eliminated 100% of citation-rendering failures by rewriting the BibTeX parser",
    "4 production-ready PRs in 2 weeks as sole student contributor",
    "Zero regressions in a six-person research codebase",
    "Tracking 200+ autoformalization papers (GAMA/LELMA)",
    "TypeScript · BibTeX · Git/GitHub · Neurosymbolic AI",
], Inches(8.3), Inches(1.88), Inches(4.8), Inches(2.2),
   size=12, color=GREY)

# Other roles
roles = [
    ("Student Representative", "BSc CS (AI), RHUL · Jan 2026–Present",
     "Elected cohort rep for 27 students. Secured flexible timetable.\nLeadership", INDIGO),
    ("Technology Officer", "Neuroscience Society · Mar 2026–Present",
     "Automated 3 events with Power Automate. Saved ~3 hrs/event.\nTech", GOLD),
]
for ri, (pos, comp, desc, col) in enumerate(roles):
    ry = Inches(4.15) + ri * Inches(1.25)
    rect(sl, Inches(8.3), ry, Inches(4.8), Inches(1.12),
         fill=CARD_BG, border=col, bw=1)
    tb(sl, pos, Inches(8.45), ry + Inches(0.08), Inches(4.5), Inches(0.28),
       size=11, bold=True, color=col)
    tb(sl, comp, Inches(8.45), ry + Inches(0.35), Inches(4.5), Inches(0.22),
       size=9, color=DIMGREY)
    tb(sl, desc, Inches(8.45), ry + Inches(0.58), Inches(4.5), Inches(0.45),
       size=9.5, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 – PAGES WALKTHROUGH (visual wireframes)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "Pages Walkthrough — Visual Overview")

page_data = [
    ("Home", PURPLE, [
        ("hero", PURPLE), ("stats", DIMGREY), ("projects", GREEN),
        ("skills", INDIGO), ("contact", ORANGE)
    ]),
    ("About", GREEN, [
        ("hero", GREEN), ("bio col", WHITE), ("edu card", INDIGO),
        ("cta", PURPLE)
    ]),
    ("Projects", GOLD, [
        ("filter tabs", GOLD), ("MyanPay", PURPLE), ("7 cards grid", INDIGO),
        ("contact footer", ORANGE)
    ]),
    ("Experience", INDIGO, [
        ("hero", INDIGO), ("DICE Lab", PURPLE), ("student rep", GREEN),
        ("tech officer", GOLD)
    ]),
    ("Contact", ORANGE, [
        ("hero", ORANGE), ("form", RED), ("direct links", GREEN),
        ("availability", DIMGREY)
    ]),
]

pcard_w = Inches(2.4)
pcard_h = Inches(5.3)
pgap    = Inches(0.22)
total_w = len(page_data) * pcard_w + (len(page_data)-1) * pgap
plx     = (W - total_w) / 2
pty     = Inches(1.2)

for pi, (pname, pcol, sections) in enumerate(page_data):
    px = plx + pi * (pcard_w + pgap)

    # Browser frame mini
    rect(sl, px, pty, pcard_w, pcard_h, fill=CODE_BG, border=pcol, bw=1.5)
    rect(sl, px, pty, pcard_w, Inches(0.28), fill=DEEP, border=None)
    ds2 = Inches(0.065)
    for di, dc in enumerate([DOT_R, DOT_Y, DOT_G]):
        circle(sl, px + Inches(0.1) + di * Inches(0.11), pty + Inches(0.14), ds2/2, dc)

    # Nav bar inside
    rect(sl, px, pty + Inches(0.28), pcard_w, Inches(0.22),
         fill=DEEP, border=None)

    # Content sections
    sec_t = pty + Inches(0.55)
    for si2, (sec_name, scol) in enumerate(sections):
        sh2 = (pcard_h - Inches(0.55) - Inches(0.2)) / len(sections) - Inches(0.06)
        sy  = sec_t + si2 * (sh2 + Inches(0.06))
        rect(sl, px + Inches(0.08), sy, pcard_w - Inches(0.16), sh2,
             fill=CARD_BG, border=scol, bw=0.75)
        tb(sl, sec_name, px + Inches(0.12), sy + sh2/2 - Pt(7),
           pcard_w - Inches(0.24), Pt(14), size=8, color=scol,
           align=PP_ALIGN.CENTER)

    # Page title below
    tb(sl, pname, px, pty + pcard_h + Inches(0.06), pcard_w, Inches(0.3),
       size=12, bold=True, color=pcol, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 – EXTERNAL REFERENCES
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)
section_bar(sl, "External References & Inspirations")

refs = [
    (PURPLE, "Google Fonts — Syne",
     "fonts.google.com/specimen/Syne",
     "Bold headings — geometric, modern feel for skill categories"),
    (INDIGO, "Google Fonts — DM Mono",
     "fonts.google.com/specimen/DM+Mono",
     "Monospace labels — gives a technical developer aesthetic"),
    (GREEN,  "MDN — CSS Grid Layout",
     "developer.mozilla.org/en-US/docs/Web/CSS/CSS_grid_layout",
     "Reference for project grid, skills grid, about page layout, video feature"),
    (GOLD,   "MDN — IntersectionObserver API",
     "developer.mozilla.org/en-US/docs/Web/API/Intersection_Observer_API",
     "Core API powering all scroll-triggered fade-in animations"),
    (ORANGE, "MDN — position: sticky",
     "developer.mozilla.org/en-US/docs/Web/CSS/position",
     "Sticky navbar: stays visible while user scrolls the page"),
    (RED,    "CSS-Tricks — A Complete Guide to Flexbox",
     "css-tricks.com/snippets/css/a-guide-to-flexbox/",
     "Reference while building the navigation bar and button rows"),
]
ref_w = Inches(5.85)
ref_h = Inches(0.92)
ref_gap = Inches(0.14)
ref_l  = [Inches(0.55), Inches(6.85)]
ref_t  = Inches(1.2)
for ri2, (col, title, url, desc) in enumerate(refs):
    rxi = ri2 % 2
    ryi = ri2 // 2
    lx  = ref_l[rxi]
    ty  = ref_t + ryi * (ref_h + ref_gap)
    rect(sl, lx, ty, ref_w, ref_h, fill=CARD_BG, border=col, bw=1.2)
    # colour accent bar left
    rect(sl, lx, ty, Inches(0.06), ref_h, fill=col, border=None)
    tb(sl, title, lx + Inches(0.15), ty + Inches(0.08),
       ref_w - Inches(0.2), Inches(0.28), size=12, bold=True, color=col)
    tb(sl, url, lx + Inches(0.15), ty + Inches(0.36),
       ref_w - Inches(0.2), Inches(0.22), size=8.5, color=DIMGREY,
       font="Courier New")
    tb(sl, desc, lx + Inches(0.15), ty + Inches(0.58),
       ref_w - Inches(0.2), Inches(0.28), size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 – SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, Inches(0.08), H, fill=PURPLE, border=None)

# Profile photo
prof_sz = Inches(2.2)
sl.shapes.add_picture(CIRCLE_PROFILE,
                       W - prof_sz - Inches(0.6),
                       Inches(0.45),
                       prof_sz, prof_sz)

section_bar(sl, "Summary")

# Left: what was built
tb(sl, "What was built",
   Inches(0.55), Inches(1.2), Inches(7.5), Inches(0.35),
   size=14, bold=True, color=PURPLE)
bullets(sl, [
    "5-page personal portfolio website — hand-coded from scratch",
    "1 shared CSS stylesheet — 1,100+ lines, zero frameworks",
    "6 JavaScript files — 1 shared module + 5 page-specific",
    "Responsive design — 600px & 700px breakpoints",
    "Dark theme with #867bff purple accent design system",
], Inches(0.55), Inches(1.6), Inches(7.5), Inches(2.0), size=13)

# Right: achievement pills
tb(sl, "Key JavaScript features",
   Inches(0.55), Inches(3.8), Inches(7.5), Inches(0.35),
   size=14, bold=True, color=PURPLE)
features2 = [
    ("IntersectionObserver scroll animations", PURPLE),
    ("Project filter tabs (data-* + classList)", GREEN),
    ("Form validation with inline errors", RED),
    ("ES module import/export pattern", INDIGO),
    ("CSS :hover dropdown (no JS)", GOLD),
    ("Video graceful fallback", ORANGE),
]
fx3 = Inches(0.55)
fy3 = Inches(4.22)
for fi3, (lbl3, col3) in enumerate(features2):
    fw3 = Inches(len(lbl3) * 0.098 + 0.55)
    if fx3 + fw3 > Inches(8.3):
        fx3 = Inches(0.55)
        fy3 += Inches(0.45)
    pill_shape(sl, lbl3, fx3, fy3, fw3, Inches(0.32),
               bg=RGBColor(0x13, 0x12, 0x1F), fg=col3)
    fx3 += fw3 + Inches(0.1)

# Key CSS features
tb(sl, "Key CSS features",
   Inches(0.55), Inches(5.42), Inches(7.5), Inches(0.35),
   size=14, bold=True, color=PURPLE)
css_feats = [
    ("CSS Grid layouts", PURPLE), ("Flexbox navbar", GREEN),
    ("Card hover transitions", INDIGO), ("Responsive @media", GOLD),
    ("CSS :hover dropdown", ORANGE), ("Animate-on-scroll classes", RED),
]
fx4 = Inches(0.55); fy4 = Inches(5.82)
for lbl4, col4 in css_feats:
    fw4 = Inches(len(lbl4) * 0.098 + 0.55)
    if fx4 + fw4 > Inches(8.3):
        fx4 = Inches(0.55); fy4 += Inches(0.45)
    pill_shape(sl, lbl4, fx4, fy4, fw4, Inches(0.32),
               bg=RGBColor(0x13, 0x12, 0x1F), fg=col4)
    fx4 += fw4 + Inches(0.1)

# Final tag
hline(sl, Inches(0.55), H - Inches(0.75), Inches(12.2), color=BORDER)
tb(sl, "HTML · CSS · JavaScript — built from scratch, no frameworks or templates",
   Inches(0.55), H - Inches(0.65), Inches(12.2), Inches(0.45),
   size=13, color=PURPLE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────
out = r"c:\Websites\MyCV\CV_Website_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")

# Cleanup temp files
for f in _tmp_files:
    try: os.remove(f)
    except: pass
